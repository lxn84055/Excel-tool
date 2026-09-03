import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import pandas as pd
import os
import threading
from datetime import datetime
import re
import csv
import time
import numpy as np

class DataProcessingGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("Excel/Word 数据处理工具")
        self.root.geometry("960x540")
        
        # 设置样式
        style = ttk.Style()
        style.theme_use('clam')
        
        # 创建主Canvas和滚动条
        self.main_canvas = tk.Canvas(root, highlightthickness=0)
        self.main_scrollbar = ttk.Scrollbar(root, orient="vertical", command=self.main_canvas.yview)
        self.main_canvas.configure(yscrollcommand=self.main_scrollbar.set)
        
        # 布局主Canvas和滚动条
        self.main_canvas.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        self.main_scrollbar.grid(row=0, column=1, sticky=(tk.N, tk.S))
        
        # 创建主框架（放在Canvas中）
        self.main_frame = ttk.Frame(self.main_canvas, padding="10")
        self.canvas_window = self.main_canvas.create_window((0, 0), window=self.main_frame, anchor="nw")
        
        # 绑定配置事件
        self.main_frame.bind("<Configure>", self.on_frame_configure)
        self.main_canvas.bind("<Configure>", self.on_canvas_configure)
        
        # 绑定鼠标滚轮事件
        self.bind_mousewheel_recursive(root)
        
        # 创建标签页
        self.notebook = ttk.Notebook(self.main_frame)
        
        # 创建各个功能标签页
        self.create_merge_tab()
        self.create_convert_tab()
        self.create_split_tab()
        self.create_batch_tab()
        
        # 创建进度条区域
        self.create_progress_area()
        
        # 日志区域
        self.create_log_area()
        
        # 初始布局：标签页在上，进度和日志在下
        self.notebook.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        self.progress_frame.grid(row=1, column=0, sticky=(tk.W, tk.E), pady=5)
        self.log_frame.grid(row=2, column=0, sticky=(tk.W, tk.E), pady=10)
        
        # 配置网格权重
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        self.main_frame.columnconfigure(0, weight=1)
        self.main_frame.rowconfigure(0, weight=1)
        
        # 初始化处理状态
        self.is_processing = False
        self.last_progress_time = 0
        
        # 设置窗口最小大小
        self.root.minsize(800, 500)
        
        # 设置窗口居中
        self.center_window()
        
        # 保存所有可交互的控件
        self.interactive_widgets = []
        self.collect_interactive_widgets()
    
    def collect_interactive_widgets(self):
        """收集所有可交互的控件"""
        for widget in self.get_all_widgets(self.main_frame):
            if isinstance(widget, (ttk.Button, ttk.Entry, ttk.Checkbutton, 
                                  ttk.Radiobutton, ttk.Combobox, tk.Listbox,
                                  ttk.Scale, ttk.Spinbox)):
                if widget != self.cancel_button:
                    self.interactive_widgets.append(widget)
    
    def lock_interface(self):
        """锁定界面"""
        self.is_processing = True
        
        for tab_id in self.notebook.tabs():
            self.notebook.tab(tab_id, state='disabled')
        
        for widget in self.interactive_widgets:
            try:
                widget.config(state='disabled')
            except:
                pass
        
        try:
            self.main_scrollbar.grid_remove()
        except:
            pass
        
        self.unbind_mousewheel_recursive(self.root)
        self.cancel_button.config(state='normal')
        
        # 处理中布局：进度和日志置顶
        self.notebook.grid_remove()
        self.progress_frame.grid_remove()
        self.log_frame.grid_remove()
        
        self.progress_frame.grid(row=0, column=0, sticky=(tk.W, tk.E), pady=5)
        self.log_frame.grid(row=1, column=0, sticky=(tk.W, tk.E), pady=5)
        self.notebook.grid(row=2, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        self.root.update_idletasks()
    
    def unlock_interface(self):
        """解锁界面"""
        self.is_processing = False
        
        for tab_id in self.notebook.tabs():
            self.notebook.tab(tab_id, state='normal')
        
        for widget in self.interactive_widgets:
            try:
                widget.config(state='normal')
            except:
                pass
        
        try:
            self.main_scrollbar.grid()
        except:
            pass
        
        self.bind_mousewheel_recursive(self.root)
        self.cancel_button.config(state='disabled')
        
        # 恢复正常布局：标签页在上，进度和日志在下
        self.notebook.grid_remove()
        self.progress_frame.grid_remove()
        self.log_frame.grid_remove()
        
        self.notebook.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        self.progress_frame.grid(row=1, column=0, sticky=(tk.W, tk.E), pady=5)
        self.log_frame.grid(row=2, column=0, sticky=(tk.W, tk.E), pady=10)
        
        self.root.update_idletasks()
    
    def unbind_mousewheel_recursive(self, widget):
        """递归解绑鼠标滚轮事件"""
        try:
            widget.unbind("<MouseWheel>")
            widget.unbind("<Button-4>")
            widget.unbind("<Button-5>")
        except:
            pass
        
        for child in widget.winfo_children():
            self.unbind_mousewheel_recursive(child)
    
    def get_all_widgets(self, parent):
        """递归获取所有子组件"""
        widgets = []
        try:
            for child in parent.winfo_children():
                widgets.append(child)
                widgets.extend(self.get_all_widgets(child))
        except:
            pass
        return widgets
    
    def detect_header(self, df):
        """检测DataFrame是否有列名（表头）"""
        if df.empty:
            return False
        
        column_names = df.columns.tolist()
        
        header_score = 0
        data_score = 0
        
        for col in column_names:
            col_str = str(col)
            if re.match(r'^[A-Za-z\u4e00-\u9fff][A-Za-z0-9\u4e00-\u9fff_\s]*$', col_str):
                header_score += 1
            if re.match(r'^\d+$', col_str):
                data_score += 1
            elif 'Unnamed' in col_str:
                data_score += 1
        
        return header_score > data_score
    
    def detect_header_quick(self, file_path):
        """快速检测文件是否有列名（只读取前几行）"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.csv':
                return self.detect_header_csv_quick(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return self.detect_header_excel_quick(file_path)
            else:
                return True
        except Exception as e:
            return True
    
    def detect_header_csv_quick(self, file_path):
        """快速检测CSV文件是否有列名"""
        try:
            encoding = self.detect_file_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                lines = []
                for i, line in enumerate(f):
                    if i >= 3:
                        break
                    lines.append(line.strip())
            
            if not lines:
                return False
            
            delimiter = self.detect_delimiter(lines[0])
            
            first_row = lines[0].split(delimiter)
            second_row = lines[1].split(delimiter) if len(lines) > 1 else None
            
            header_score = 0
            data_score = 0
            
            for cell in first_row:
                cell = cell.strip().strip('"').strip("'")
                if re.match(r'^[A-Za-z\u4e00-\u9fff][A-Za-z0-9\u4e00-\u9fff_\s]*$', cell):
                    header_score += 1
                if re.match(r'^\d+$', cell):
                    data_score += 1
            
            if second_row:
                for cell1, cell2 in zip(first_row, second_row):
                    cell1 = cell1.strip().strip('"').strip("'")
                    cell2 = cell2.strip().strip('"').strip("'")
                    if re.match(r'^[A-Za-z\u4e00-\u9fff]', cell1) and re.match(r'^\d+$', cell2):
                        header_score += 2
            
            return header_score > data_score
            
        except Exception as e:
            return True
    
    def detect_header_excel_quick(self, file_path):
        """快速检测Excel文件是否有列名（只读取前几行）"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.xlsx':
                df = pd.read_excel(file_path, engine='openpyxl', nrows=3, dtype=object)
            else:
                try:
                    df = pd.read_excel(file_path, engine='xlrd', nrows=3, dtype=object)
                except:
                    df = pd.read_excel(file_path, engine='openpyxl', nrows=3, dtype=object)
            
            if df.empty:
                return False
            
            column_names = df.columns.tolist()
            
            header_score = 0
            data_score = 0
            
            for col in column_names:
                col_str = str(col)
                if re.match(r'^[A-Za-z\u4e00-\u9fff][A-Za-z0-9\u4e00-\u9fff_\s]*$', col_str):
                    header_score += 1
                if re.match(r'^\d+$', col_str):
                    data_score += 1
                elif 'Unnamed' in col_str:
                    data_score += 1
            
            return header_score > data_score
            
        except Exception as e:
            return True
    
    def get_columns_quick(self, file_path):
        """快速获取列名（只读取表头）"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            has_header = self.detect_header_quick(file_path)
            
            if file_extension == '.csv':
                return self.get_csv_columns_quick(file_path, has_header)
            elif file_extension in ['.xlsx', '.xls']:
                return self.get_excel_columns_quick(file_path, has_header)
            else:
                return [], True
        except Exception as e:
            return [], True
    
    def get_csv_columns_quick(self, file_path, has_header):
        """快速获取CSV列名"""
        try:
            encoding = self.detect_file_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                lines = []
                for i, line in enumerate(f):
                    if i >= 5:
                        break
                    lines.append(line.strip())
            
            if not lines:
                return [], has_header
            
            delimiter = self.detect_delimiter(lines[0])
            
            if has_header:
                columns = [c.strip().strip('"').strip("'") for c in lines[0].split(delimiter)]
            else:
                max_cols = max(len(line.split(delimiter)) for line in lines)
                columns = [f"列{i+1}" for i in range(max_cols)]
            
            return columns, has_header
            
        except Exception as e:
            return [], has_header
    
    def get_excel_columns_quick(self, file_path, has_header):
        """快速获取Excel列名"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.xlsx':
                df = pd.read_excel(file_path, engine='openpyxl', nrows=5, dtype=object)
            else:
                try:
                    df = pd.read_excel(file_path, engine='xlrd', nrows=5, dtype=object)
                except:
                    df = pd.read_excel(file_path, engine='openpyxl', nrows=5, dtype=object)
            
            if has_header:
                columns = df.columns.tolist()
            else:
                max_cols = len(df.columns)
                columns = [f"列{i+1}" for i in range(max_cols)]
            
            return columns, has_header
            
        except Exception as e:
            return [], has_header
    
    def normalize_columns(self, df, reference_columns=None):
        """统一列名"""
        if reference_columns is None:
            reference_columns = df.columns.tolist()
        
        current_columns = df.columns.tolist()
        
        if len(current_columns) == len(reference_columns):
            if current_columns != reference_columns:
                df.columns = reference_columns
        else:
            matched_columns = []
            for i in range(min(len(current_columns), len(reference_columns))):
                matched_columns.append(reference_columns[i])
            
            while len(matched_columns) < len(current_columns):
                matched_columns.append(f"列{len(matched_columns)+1}")
            
            df.columns = matched_columns
        
        return df
    
    def read_csv_without_header(self, file_path):
        """读取无列名的CSV文件"""
        encoding = self.detect_file_encoding(file_path)
        
        with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
            sample = f.read(4096)
            f.seek(0)
            delimiter = self.detect_delimiter(sample)
            
            csv_reader = csv.reader(f, delimiter=delimiter)
            rows = [row for row in csv_reader if row]
        
        if not rows:
            return pd.DataFrame()
        
        max_cols = max(len(row) for row in rows)
        default_columns = [f"列{i+1}" for i in range(max_cols)]
        
        processed_rows = []
        for row in rows:
            if len(row) < max_cols:
                row = row + [''] * (max_cols - len(row))
            elif len(row) > max_cols:
                row = row[:max_cols]
            processed_rows.append(row)
        
        df = pd.DataFrame(processed_rows, columns=default_columns, dtype=object)
        return df
    
    def read_excel_without_header(self, file_path):
        """读取无列名的Excel文件"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.xlsx':
            df = pd.read_excel(file_path, engine='openpyxl', header=None, dtype=object)
        else:
            try:
                df = pd.read_excel(file_path, engine='xlrd', header=None, dtype=object)
            except:
                df = pd.read_excel(file_path, engine='openpyxl', header=None, dtype=object)
        
        default_columns = [f"列{i+1}" for i in range(len(df.columns))]
        df.columns = default_columns
        return df
    
    def remove_blank_data(self, df, remove_blank_rows=True, remove_blank_cols=True, 
                          remove_blank_cells=True):
        """去除空白数据"""
        original_shape = df.shape
        
        if remove_blank_rows:
            df = df.dropna(how='all')
            if original_shape[0] - df.shape[0] > 0:
                self.log(f"去除空白行: {original_shape[0] - df.shape[0]}行")
        
        if remove_blank_cols:
            df = df.dropna(axis=1, how='all')
            if original_shape[1] - df.shape[1] > 0:
                self.log(f"去除空白列: {original_shape[1] - df.shape[1]}列")
        
        if remove_blank_cells:
            for col in df.columns:
                if df[col].dtype == 'object':
                    df[col] = df[col].apply(lambda x: x.strip() if isinstance(x, str) else x)
                    df[col] = df[col].apply(lambda x: np.nan if isinstance(x, str) and x == '' else x)
        
        return df
    
    def convert_dtypes_after_processing(self, df):
        """处理后将文本形式的数字转换为数字类型，但保留日期时间"""
        try:
            for col in df.columns:
                if df[col].dtype in ['int64', 'float64', 'datetime64[ns]', 'Int64', 'Float64']:
                    continue
                
                if self.is_datetime_column(df[col]):
                    continue
                
                try:
                    numeric_values = pd.to_numeric(df[col], errors='coerce')
                    non_null_count = df[col].notna().sum()
                    numeric_count = numeric_values.notna().sum()
                    
                    if numeric_count > non_null_count * 0.8:
                        if (numeric_values.dropna() == numeric_values.dropna().astype(int)).all():
                            df[col] = numeric_values.astype('Int64')
                        else:
                            df[col] = numeric_values
                except:
                    pass
            
            return df
        except Exception as e:
            return df
    
    def is_datetime_column(self, series):
        """检查列是否包含日期时间数据"""
        try:
            sample = series.dropna().head(100)
            if len(sample) == 0:
                return False
            
            if pd.api.types.is_datetime64_any_dtype(series):
                return True
            
            datetime_count = 0
            for value in sample:
                if isinstance(value, (datetime, pd.Timestamp)):
                    datetime_count += 1
                elif isinstance(value, str):
                    date_patterns = [
                        r'\d{4}-\d{2}-\d{2}',
                        r'\d{4}/\d{2}/\d{2}',
                        r'\d{2}-\d{2}-\d{4}',
                        r'\d{2}/\d{2}/\d{4}',
                        r'\d{4}年\d{1,2}月\d{1,2}日',
                        r'\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}',
                        r'\d{4}/\d{2}/\d{2} \d{2}:\d{2}:\d{2}',
                    ]
                    for pattern in date_patterns:
                        if re.search(pattern, value):
                            datetime_count += 1
                            break
            
            return datetime_count > len(sample) * 0.5
        except:
            return False
    
    def read_csv_file_robust(self, file_path):
        """读取CSV文件"""
        try:
            encoding = self.detect_file_encoding(file_path)
            
            rows = []
            max_columns = 0
            
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                sample = f.read(4096)
                f.seek(0)
                delimiter = self.detect_delimiter(sample)
                
                csv_reader = csv.reader(f, delimiter=delimiter)
                
                for row in csv_reader:
                    if row:
                        rows.append(row)
                        max_columns = max(max_columns, len(row))
            
            if not rows:
                raise Exception("CSV文件为空")
            
            processed_rows = []
            for row in rows:
                if len(row) < max_columns:
                    row = row + [''] * (max_columns - len(row))
                elif len(row) > max_columns:
                    row = row[:max_columns]
                processed_rows.append(row)
            
            if len(processed_rows) > 1:
                columns = processed_rows[0]
                data = processed_rows[1:]
                df = pd.DataFrame(data, columns=columns, dtype=object)
            else:
                df = pd.DataFrame(processed_rows, dtype=object)
            
            return df
            
        except Exception as e:
            raise e
    
    def detect_file_encoding(self, file_path):
        """检测文件编码"""
        encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    f.read(1024)
                return encoding
            except:
                continue
        
        return 'latin1'
    
    def detect_delimiter(self, sample_text):
        """检测CSV分隔符"""
        delimiters = [',', ';', '\t', '|']
        
        counts = {}
        for delimiter in delimiters:
            counts[delimiter] = sample_text.count(delimiter)
        
        if counts:
            max_delimiter = max(counts, key=counts.get)
            if counts[max_delimiter] > 0:
                return max_delimiter
        
        return ','
    
    def read_excel_file(self, file_path):
        """智能读取文件"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.csv':
                return self.read_csv_file_robust(file_path)
            elif file_extension == '.xlsx':
                return pd.read_excel(file_path, engine='openpyxl', dtype=object)
            elif file_extension == '.xls':
                try:
                    return pd.read_excel(file_path, engine='xlrd', dtype=object)
                except:
                    return pd.read_excel(file_path, engine='openpyxl', dtype=object)
            else:
                return pd.read_excel(file_path, dtype=object)
        
        except Exception as e:
            raise e
    
    def save_excel_file(self, df, file_path):
        """保存文件，应用数据类型转换"""
        try:
            df = self.convert_dtypes_after_processing(df)
            
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.csv':
                df.to_csv(file_path, index=False, encoding='utf-8-sig')
                return True
            elif file_extension == '.xlsx':
                df.to_excel(file_path, index=False, engine='openpyxl')
                return True
            elif file_extension == '.xls':
                df.to_excel(file_path, index=False, engine='xlwt')
                return True
            else:
                df.to_excel(file_path, index=False, engine='openpyxl')
                return True
        
        except Exception as e:
            return False
    
    def on_frame_configure(self, event):
        self.main_canvas.configure(scrollregion=self.main_canvas.bbox("all"))
    
    def on_canvas_configure(self, event):
        self.main_canvas.itemconfig(self.canvas_window, width=event.width)
    
    def bind_mousewheel_recursive(self, widget):
        """绑定鼠标滚轮事件"""
        widget.bind("<MouseWheel>", self.on_mousewheel_windows)
        widget.bind("<Button-4>", self.on_mousewheel_linux)
        widget.bind("<Button-5>", self.on_mousewheel_linux)
        
        for child in widget.winfo_children():
            self.bind_mousewheel_recursive(child)
    
    def on_mousewheel_windows(self, event):
        if self.is_processing:
            return "break"
        
        if event.delta > 0:
            self.main_canvas.yview_scroll(-3, "units")
        else:
            self.main_canvas.yview_scroll(3, "units")
        return "break"
    
    def on_mousewheel_linux(self, event):
        if self.is_processing:
            return "break"
        
        if event.num == 4:
            self.main_canvas.yview_scroll(-3, "units")
        elif event.num == 5:
            self.main_canvas.yview_scroll(3, "units")
        return "break"
    
    def center_window(self):
        self.root.update_idletasks()
        width = 960
        height = 540
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f'{width}x{height}+{x}+{y}')
    
    def create_progress_area(self):
        self.progress_frame = ttk.LabelFrame(self.main_frame, text="处理进度", padding="5")
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(self.progress_frame, variable=self.progress_var, 
                                           maximum=100, length=900, mode='determinate')
        self.progress_bar.grid(row=0, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        
        self.progress_label = ttk.Label(self.progress_frame, text="0%")
        self.progress_label.grid(row=0, column=2, padx=10)
        
        self.status_label = ttk.Label(self.progress_frame, text="就绪", foreground="green")
        self.status_label.grid(row=1, column=0, columnspan=3, sticky=tk.W, pady=5)
        
        self.indeterminate_progress = ttk.Progressbar(self.progress_frame, mode='indeterminate',
                                                      length=900)
        self.indeterminate_progress.grid(row=2, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        self.indeterminate_progress.grid_remove()
        
        self.cancel_button = ttk.Button(self.progress_frame, text="取消", command=self.cancel_operation,
                                       state='disabled')
        self.cancel_button.grid(row=2, column=2, padx=10)
        
        self.cancel_flag = False
        self.last_progress_time = 0
    
    def create_log_area(self):
        self.log_frame = ttk.LabelFrame(self.main_frame, text="操作日志", padding="5")
        
        self.log_text = scrolledtext.ScrolledText(self.log_frame, height=4, width=100)
        self.log_text.grid(row=0, column=0, sticky=(tk.W, tk.E))
        
        clear_btn = ttk.Button(self.log_frame, text="清除日志", command=self.clear_log)
        clear_btn.grid(row=1, column=0, pady=2)
    
    def clear_log(self):
        if self.is_processing:
            return
        self.log_text.delete(1.0, tk.END)
    
    def log(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
        self.root.update_idletasks()
    
    def update_progress(self, value, status_text=None, force_update=False):
        current_time = time.time()
        if not force_update and current_time - self.last_progress_time < 0.1:
            return
        
        self.last_progress_time = current_time
        self.progress_var.set(value)
        self.progress_label.config(text=f"{int(value)}%")
        if status_text:
            self.status_label.config(text=status_text)
        self.root.update_idletasks()
    
    def start_indeterminate(self, status_text="处理中..."):
        self.indeterminate_progress.grid()
        self.indeterminate_progress.start(10)
        self.status_label.config(text=status_text, foreground="blue")
        self.cancel_button.config(state='normal')
        self.cancel_flag = False
        
        self.lock_interface()
        
        self.root.update_idletasks()
    
    def stop_indeterminate(self, status_text="完成", success=True):
        self.indeterminate_progress.stop()
        self.indeterminate_progress.grid_remove()
        if success:
            self.status_label.config(text=status_text, foreground="green")
            self.progress_var.set(100)
            self.progress_label.config(text="100%")
        else:
            self.status_label.config(text=status_text, foreground="red")
        
        self.unlock_interface()
        
        self.root.update_idletasks()
    
    def cancel_operation(self):
        self.cancel_flag = True
        self.status_label.config(text="正在取消...", foreground="orange")
        self.cancel_button.config(state='disabled')
        self.log("用户请求取消操作")
    
    def check_cancel(self):
        if self.cancel_flag:
            raise Exception("操作已被用户取消")
    
    def create_merge_tab(self):
        """创建数据合并标签页"""
        merge_frame = ttk.Frame(self.notebook, padding="10")
        self.notebook.add(merge_frame, text="数据合并")
        
        file_frame = ttk.LabelFrame(merge_frame, text="选择要合并的文件（支持Excel和CSV）", padding="10")
        file_frame.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        file_list_frame = ttk.Frame(file_frame)
        file_list_frame.grid(row=0, column=0, columnspan=3, pady=5, sticky=(tk.W, tk.E))
        
        self.file_listbox = tk.Listbox(file_list_frame, height=3, width=70, selectmode=tk.MULTIPLE)
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        file_scrollbar = ttk.Scrollbar(file_list_frame, orient="vertical", command=self.file_listbox.yview)
        file_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.file_listbox.configure(yscrollcommand=file_scrollbar.set)
        
        button_frame = ttk.Frame(file_frame)
        button_frame.grid(row=1, column=0, columnspan=3, pady=5)
        
        ttk.Button(button_frame, text="添加文件", command=self.add_files).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="移除选中", command=self.remove_selected_files).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="清空列表", command=self.clear_file_list).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="预览列名", command=self.preview_columns).pack(side=tk.LEFT, padx=5)
        
        column_frame = ttk.LabelFrame(merge_frame, text="列选择（留空表示选择所有列）", padding="10")
        column_frame.grid(row=1, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(column_frame, text="按列名选择:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.columns_by_name = ttk.Entry(column_frame, width=50)
        self.columns_by_name.grid(row=0, column=1, columnspan=2, sticky=tk.W, pady=3)
        ttk.Label(column_frame, text="(用逗号分隔列名)").grid(row=0, column=3, sticky=tk.W)
        
        ttk.Label(column_frame, text="按列坐标选择:").grid(row=1, column=0, sticky=tk.W, pady=3)
        self.columns_by_index = ttk.Entry(column_frame, width=50)
        self.columns_by_index.grid(row=1, column=1, columnspan=2, sticky=tk.W, pady=3)
        ttk.Label(column_frame, text="(如: 1,3,5 或 1-5)").grid(row=1, column=3, sticky=tk.W)
        
        option_frame = ttk.LabelFrame(merge_frame, text="合并选项", padding="10")
        option_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(option_frame, text="合并方式:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.merge_type = tk.StringVar(value="vertical")
        ttk.Radiobutton(option_frame, text="垂直合并", variable=self.merge_type, 
                       value="vertical").grid(row=0, column=1, sticky=tk.W, padx=5)
        ttk.Radiobutton(option_frame, text="水平合并", variable=self.merge_type, 
                       value="horizontal").grid(row=0, column=2, sticky=tk.W, padx=5)
        
        self.remove_blank_rows_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="去空白行", 
                       variable=self.remove_blank_rows_var).grid(row=1, column=0, sticky=tk.W, pady=3)
        
        self.remove_blank_cols_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="去空白列", 
                       variable=self.remove_blank_cols_var).grid(row=1, column=1, sticky=tk.W, pady=3)
        
        self.remove_blank_cells_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="去单元格空白", 
                       variable=self.remove_blank_cells_var).grid(row=1, column=2, sticky=tk.W, pady=3)
        
        self.add_source_var = tk.BooleanVar()
        ttk.Checkbutton(option_frame, text="添加来源标识", 
                       variable=self.add_source_var).grid(row=2, column=0, sticky=tk.W, pady=3)
        
        self.remove_dup_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="去除重复行", 
                       variable=self.remove_dup_var).grid(row=2, column=1, sticky=tk.W, pady=3)
        
        output_frame = ttk.LabelFrame(merge_frame, text="输出设置", padding="10")
        output_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(output_frame, text="输出文件:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.output_path = tk.StringVar()
        ttk.Entry(output_frame, textvariable=self.output_path, width=50).grid(row=0, column=1, padx=5)
        ttk.Button(output_frame, text="浏览", command=self.select_output_file).grid(row=0, column=2)
        
        self.merge_button = ttk.Button(merge_frame, text="开始合并", command=self.start_merge, width=15)
        self.merge_button.grid(row=4, column=0, columnspan=3, pady=5)
    
    def create_convert_tab(self):
        """创建格式转换标签页"""
        convert_frame = ttk.Frame(self.notebook, padding="10")
        self.notebook.add(convert_frame, text="格式转换")
        
        type_frame = ttk.LabelFrame(convert_frame, text="转换类型", padding="10")
        type_frame.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.convert_type = tk.StringVar(value="excel_to_word")
        
        ttk.Radiobutton(type_frame, text="Excel→Word", variable=self.convert_type, 
                       value="excel_to_word").grid(row=0, column=0, padx=5, pady=3)
        ttk.Radiobutton(type_frame, text="Word→Excel", variable=self.convert_type, 
                       value="word_to_excel").grid(row=0, column=1, padx=5, pady=3)
        ttk.Radiobutton(type_frame, text="Excel→PDF", variable=self.convert_type, 
                       value="excel_to_pdf").grid(row=0, column=2, padx=5, pady=3)
        ttk.Radiobutton(type_frame, text="Excel→PPT", variable=self.convert_type, 
                       value="excel_to_ppt").grid(row=1, column=0, padx=5, pady=3)
        ttk.Radiobutton(type_frame, text="Word→PDF", variable=self.convert_type, 
                       value="word_to_pdf").grid(row=1, column=1, padx=5, pady=3)
        ttk.Radiobutton(type_frame, text="PPT→PDF", variable=self.convert_type, 
                       value="ppt_to_pdf").grid(row=1, column=2, padx=5, pady=3)
        
        input_frame = ttk.LabelFrame(convert_frame, text="输入文件", padding="10")
        input_frame.grid(row=1, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.convert_input = tk.StringVar()
        ttk.Entry(input_frame, textvariable=self.convert_input, width=60).grid(row=0, column=0, padx=5)
        ttk.Button(input_frame, text="浏览", command=self.select_input_file).grid(row=0, column=1)
        
        output_frame = ttk.LabelFrame(convert_frame, text="输出文件", padding="10")
        output_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.convert_output = tk.StringVar()
        ttk.Entry(output_frame, textvariable=self.convert_output, width=60).grid(row=0, column=0, padx=5)
        ttk.Button(output_frame, text="浏览", command=self.select_convert_output).grid(row=0, column=1)
        
        self.convert_button = ttk.Button(convert_frame, text="开始转换", command=self.start_convert, width=15)
        self.convert_button.grid(row=3, column=0, columnspan=3, pady=5)
    
    def create_split_tab(self):
        """创建数据拆分标签页"""
        split_frame = ttk.Frame(self.notebook, padding="10")
        self.notebook.add(split_frame, text="数据拆分")
        
        file_frame = ttk.LabelFrame(split_frame, text="选择文件", padding="10")
        file_frame.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.split_file_path = tk.StringVar()
        ttk.Entry(file_frame, textvariable=self.split_file_path, width=60).grid(row=0, column=0, padx=5)
        ttk.Button(file_frame, text="浏览", command=self.select_split_file).grid(row=0, column=1)
        ttk.Button(file_frame, text="预览列名", command=self.preview_split_columns).grid(row=0, column=2, padx=5)
        
        method_frame = ttk.LabelFrame(split_frame, text="拆分方式", padding="10")
        method_frame.grid(row=1, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.split_method = tk.StringVar(value="by_columns")
        ttk.Radiobutton(method_frame, text="按列值", variable=self.split_method, 
                       value="by_columns", command=self.toggle_split_method).grid(row=0, column=0, padx=5)
        ttk.Radiobutton(method_frame, text="按列位置", variable=self.split_method, 
                       value="by_column_position", command=self.toggle_split_method).grid(row=0, column=1, padx=5)
        ttk.Radiobutton(method_frame, text="按行数", variable=self.split_method, 
                       value="by_rows", command=self.toggle_split_method).grid(row=0, column=2, padx=5)
        ttk.Radiobutton(method_frame, text="按特定行", variable=self.split_method, 
                       value="by_specific_rows", command=self.toggle_split_method).grid(row=0, column=3, padx=5)
        
        self.column_split_frame = ttk.LabelFrame(split_frame, text="按列值拆分设置", padding="10")
        self.column_split_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(self.column_split_frame, text="拆分依据列:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.split_columns = ttk.Entry(self.column_split_frame, width=40)
        self.split_columns.grid(row=0, column=1, sticky=tk.W, pady=3)
        
        self.column_position_frame = ttk.LabelFrame(split_frame, text="按列位置拆分设置", padding="10")
        self.column_position_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        self.column_position_type = tk.StringVar(value="before_after")
        ttk.Radiobutton(self.column_position_frame, text="前后拆分", variable=self.column_position_type, 
                       value="before_after", command=self.toggle_column_position_type).grid(row=0, column=0, padx=5)
        ttk.Radiobutton(self.column_position_frame, text="指定列提取", variable=self.column_position_type, 
                       value="specific_columns", command=self.toggle_column_position_type).grid(row=0, column=1, padx=5)
        
        self.before_after_frame = ttk.Frame(self.column_position_frame)
        self.before_after_frame.grid(row=1, column=0, columnspan=3, sticky=tk.W, pady=3)
        
        ttk.Label(self.before_after_frame, text="拆分位置列号:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.split_position = ttk.Entry(self.before_after_frame, width=15)
        self.split_position.grid(row=0, column=1, sticky=tk.W, pady=3)
        ttk.Button(self.before_after_frame, text="从预览选择", command=self.preview_position_columns).grid(row=0, column=2, padx=5)
        
        self.specific_columns_frame = ttk.Frame(self.column_position_frame)
        self.specific_columns_frame.grid(row=2, column=0, columnspan=3, sticky=tk.W, pady=3)
        
        ttk.Label(self.specific_columns_frame, text="指定列号:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.specific_column_numbers = ttk.Entry(self.specific_columns_frame, width=40)
        self.specific_column_numbers.grid(row=0, column=1, sticky=tk.W, pady=3)
        ttk.Button(self.specific_columns_frame, text="从预览选择", command=self.preview_specific_columns).grid(row=0, column=2, padx=5)
        
        self.row_split_frame = ttk.LabelFrame(split_frame, text="按行数拆分设置", padding="10")
        self.row_split_frame.grid(row=4, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(self.row_split_frame, text="每个文件行数:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.rows_per_file = ttk.Entry(self.row_split_frame, width=15)
        self.rows_per_file.grid(row=0, column=1, sticky=tk.W, pady=3)
        self.rows_per_file.insert(0, "1000")
        
        self.specific_row_frame = ttk.LabelFrame(split_frame, text="按特定行拆分设置", padding="10")
        self.specific_row_frame.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(self.specific_row_frame, text="拆分行号:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.specific_rows = ttk.Entry(self.specific_row_frame, width=40)
        self.specific_rows.grid(row=0, column=1, sticky=tk.W, pady=3)
        
        output_frame = ttk.LabelFrame(split_frame, text="输出设置", padding="10")
        output_frame.grid(row=6, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(output_frame, text="输出目录:").grid(row=0, column=0, sticky=tk.W, pady=3)
        self.split_output_dir = tk.StringVar(value="./split_output")
        ttk.Entry(output_frame, textvariable=self.split_output_dir, width=40).grid(row=0, column=1, pady=3)
        ttk.Button(output_frame, text="浏览", command=self.select_split_output).grid(row=0, column=2, padx=5)
        
        self.split_button = ttk.Button(split_frame, text="开始拆分", command=self.start_split, width=15)
        self.split_button.grid(row=7, column=0, columnspan=3, pady=5)
        
        self.column_position_frame.grid_remove()
        self.row_split_frame.grid_remove()
        self.specific_row_frame.grid_remove()
        self.specific_columns_frame.grid_remove()
    
    def create_batch_tab(self):
        """创建批量处理标签页"""
        batch_frame = ttk.Frame(self.notebook, padding="10")
        self.notebook.add(batch_frame, text="批量处理")
        
        dir_frame = ttk.LabelFrame(batch_frame, text="选择文件夹", padding="10")
        dir_frame.grid(row=0, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        
        self.batch_dir = tk.StringVar()
        ttk.Entry(dir_frame, textvariable=self.batch_dir, width=60).grid(row=0, column=0, padx=5)
        ttk.Button(dir_frame, text="浏览", command=self.select_batch_dir).grid(row=0, column=1)
        
        option_frame = ttk.LabelFrame(batch_frame, text="处理选项", padding="10")
        option_frame.grid(row=1, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        
        self.batch_operation = tk.StringVar(value="clean")
        ttk.Radiobutton(option_frame, text="清理数据", variable=self.batch_operation, 
                       value="clean").grid(row=0, column=0, padx=10)
        ttk.Radiobutton(option_frame, text="去除空行", variable=self.batch_operation, 
                       value="remove_empty").grid(row=0, column=1, padx=10)
        
        self.batch_button = ttk.Button(batch_frame, text="开始批量处理", command=self.start_batch, width=15)
        self.batch_button.grid(row=2, column=0, columnspan=2, pady=5)
    
    def toggle_split_method(self):
        """切换拆分方式"""
        method = self.split_method.get()
        
        self.column_split_frame.grid_remove()
        self.column_position_frame.grid_remove()
        self.row_split_frame.grid_remove()
        self.specific_row_frame.grid_remove()
        
        if method == "by_columns":
            self.column_split_frame.grid()
        elif method == "by_column_position":
            self.column_position_frame.grid()
            self.toggle_column_position_type()
        elif method == "by_rows":
            self.row_split_frame.grid()
        elif method == "by_specific_rows":
            self.specific_row_frame.grid()
    
    def toggle_column_position_type(self):
        """切换列位置拆分类型"""
        pos_type = self.column_position_type.get()
        
        self.before_after_frame.grid_remove()
        self.specific_columns_frame.grid_remove()
        
        if pos_type == "before_after":
            self.before_after_frame.grid()
        elif pos_type == "specific_columns":
            self.specific_columns_frame.grid()
    
    def preview_columns(self):
        """预览统一后的列名"""
        if self.is_processing:
            return
        
        files = list(self.file_listbox.get(0, tk.END))
        if not files:
            messagebox.showwarning("警告", "请先添加文件")
            return
        
        try:
            reference_columns = None
            
            for file in files:
                columns, has_header = self.get_columns_quick(file)
                if has_header:
                    reference_columns = columns
                    break
            
            if reference_columns is None:
                columns, _ = self.get_columns_quick(files[0])
                reference_columns = columns if columns else [f"列{i+1}" for i in range(10)]
            
            preview_window = tk.Toplevel(self.root)
            preview_window.title("列名预览")
            preview_window.geometry("400x500")
            
            columns_listbox = tk.Listbox(preview_window, selectmode=tk.MULTIPLE, 
                                        height=20, width=50)
            columns_listbox.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)
            
            for i, col in enumerate(reference_columns, 1):
                columns_listbox.insert(tk.END, f"{i}. {col}")
            
            button_frame = ttk.Frame(preview_window)
            button_frame.pack(pady=10)
            
            def add_selected_columns():
                selected_indices = columns_listbox.curselection()
                if not selected_indices:
                    return
                
                selected_columns = []
                for idx in selected_indices:
                    display_text = columns_listbox.get(idx)
                    column_name = re.sub(r'^\d+\.\s*', '', display_text)
                    selected_columns.append(column_name)
                
                current_text = self.columns_by_name.get().strip()
                
                if current_text:
                    existing_columns = [c.strip() for c in current_text.split(',') if c.strip()]
                    for col in selected_columns:
                        if col not in existing_columns:
                            existing_columns.append(col)
                    new_text = ', '.join(existing_columns)
                else:
                    new_text = ', '.join(selected_columns)
                
                self.columns_by_name.delete(0, tk.END)
                self.columns_by_name.insert(0, new_text)
                preview_window.destroy()
            
            ttk.Button(button_frame, text="添加选中列名", 
                      command=add_selected_columns).pack(side=tk.LEFT, padx=5)
            ttk.Button(button_frame, text="关闭", 
                      command=preview_window.destroy).pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            messagebox.showerror("错误", f"读取文件失败: {str(e)}")
    
    def preview_split_columns(self):
        """预览拆分文件的列名"""
        if self.is_processing:
            return
        
        file_path = self.split_file_path.get()
        if not file_path:
            messagebox.showwarning("警告", "请先选择文件")
            return
        
        try:
            columns, has_header = self.get_columns_quick(file_path)
            
            preview_window = tk.Toplevel(self.root)
            preview_window.title("列名预览")
            preview_window.geometry("400x500")
            
            columns_listbox = tk.Listbox(preview_window, selectmode=tk.MULTIPLE, 
                                        height=20, width=50)
            columns_listbox.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)
            
            for i, col in enumerate(columns, 1):
                columns_listbox.insert(tk.END, f"{i}. {col}")
            
            button_frame = ttk.Frame(preview_window)
            button_frame.pack(pady=10)
            
            def add_selected_columns():
                selected_indices = columns_listbox.curselection()
                if not selected_indices:
                    return
                
                selected_columns = []
                for idx in selected_indices:
                    display_text = columns_listbox.get(idx)
                    column_name = re.sub(r'^\d+\.\s*', '', display_text)
                    selected_columns.append(column_name)
                
                current_text = self.split_columns.get().strip()
                
                if current_text:
                    existing_columns = [c.strip() for c in current_text.split(',') if c.strip()]
                    for col in selected_columns:
                        if col not in existing_columns:
                            existing_columns.append(col)
                    new_text = ', '.join(existing_columns)
                else:
                    new_text = ', '.join(selected_columns)
                
                self.split_columns.delete(0, tk.END)
                self.split_columns.insert(0, new_text)
                preview_window.destroy()
            
            ttk.Button(button_frame, text="添加选中列名", 
                      command=add_selected_columns).pack(side=tk.LEFT, padx=5)
            ttk.Button(button_frame, text="关闭", 
                      command=preview_window.destroy).pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            messagebox.showerror("错误", f"读取文件失败: {str(e)}")
    
    def preview_position_columns(self):
        """预览列名用于选择拆分位置"""
        if self.is_processing:
            return
        
        file_path = self.split_file_path.get()
        if not file_path:
            messagebox.showwarning("警告", "请先选择文件")
            return
        
        try:
            columns, has_header = self.get_columns_quick(file_path)
            
            preview_window = tk.Toplevel(self.root)
            preview_window.title("选择拆分位置")
            preview_window.geometry("400x500")
            
            columns_listbox = tk.Listbox(preview_window, selectmode=tk.SINGLE, 
                                        height=20, width=50)
            columns_listbox.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)
            
            for i, col in enumerate(columns, 1):
                columns_listbox.insert(tk.END, f"{i}. {col}")
            
            button_frame = ttk.Frame(preview_window)
            button_frame.pack(pady=10)
            
            def set_split_position():
                selected_indices = columns_listbox.curselection()
                if not selected_indices:
                    return
                
                selected_idx = selected_indices[0]
                column_number = selected_idx + 1
                
                self.split_position.delete(0, tk.END)
                self.split_position.insert(0, str(column_number))
                preview_window.destroy()
            
            ttk.Button(button_frame, text="设置拆分位置", 
                      command=set_split_position).pack(side=tk.LEFT, padx=5)
            ttk.Button(button_frame, text="关闭", 
                      command=preview_window.destroy).pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            messagebox.showerror("错误", f"读取文件失败: {str(e)}")
    
    def preview_specific_columns(self):
        """预览列名用于选择指定列"""
        if self.is_processing:
            return
        
        file_path = self.split_file_path.get()
        if not file_path:
            messagebox.showwarning("警告", "请先选择文件")
            return
        
        try:
            columns, has_header = self.get_columns_quick(file_path)
            
            preview_window = tk.Toplevel(self.root)
            preview_window.title("选择指定列")
            preview_window.geometry("400x500")
            
            columns_listbox = tk.Listbox(preview_window, selectmode=tk.MULTIPLE, 
                                        height=20, width=50)
            columns_listbox.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)
            
            for i, col in enumerate(columns, 1):
                columns_listbox.insert(tk.END, f"{i}. {col}")
            
            button_frame = ttk.Frame(preview_window)
            button_frame.pack(pady=10)
            
            def add_selected_columns():
                selected_indices = columns_listbox.curselection()
                if not selected_indices:
                    return
                
                selected_numbers = []
                for idx in selected_indices:
                    column_number = idx + 1
                    selected_numbers.append(str(column_number))
                
                current_text = self.specific_column_numbers.get().strip()
                
                if current_text:
                    existing_numbers = [n.strip() for n in current_text.split(',') if n.strip()]
                    for num in selected_numbers:
                        if num not in existing_numbers:
                            existing_numbers.append(num)
                    new_text = ','.join(existing_numbers)
                else:
                    new_text = ','.join(selected_numbers)
                
                self.specific_column_numbers.delete(0, tk.END)
                self.specific_column_numbers.insert(0, new_text)
                preview_window.destroy()
            
            ttk.Button(button_frame, text="添加选中列号", 
                      command=add_selected_columns).pack(side=tk.LEFT, padx=5)
            ttk.Button(button_frame, text="关闭", 
                      command=preview_window.destroy).pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            messagebox.showerror("错误", f"读取文件失败: {str(e)}")
    
    def parse_column_selection(self, columns_str, df_columns):
        """解析列选择字符串"""
        selected_columns = []
        
        if self.columns_by_name.get().strip():
            column_names = [c.strip() for c in self.columns_by_name.get().split(',') if c.strip()]
            for col_name in column_names:
                if col_name in df_columns:
                    selected_columns.append(col_name)
        
        if self.columns_by_index.get().strip():
            index_str = self.columns_by_index.get().strip()
            indices = []
            
            for part in index_str.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = part.split('-')
                    indices.extend(range(int(start), int(end) + 1))
                else:
                    indices.append(int(part))
            
            for idx in indices:
                if 1 <= idx <= len(df_columns):
                    selected_columns.append(df_columns[idx - 1])
        
        selected_columns = list(dict.fromkeys(selected_columns))
        
        return selected_columns
    
    def add_files(self):
        if self.is_processing:
            return
        
        files = filedialog.askopenfilenames(
            title="选择文件",
            filetypes=[
                ("所有支持的文件", "*.xlsx *.xls *.csv"),
                ("Excel文件", "*.xlsx *.xls"),
                ("CSV文件", "*.csv"),
                ("所有文件", "*.*")
            ]
        )
        for file in files:
            if file not in self.file_listbox.get(0, tk.END):
                self.file_listbox.insert(tk.END, file)
                self.log(f"添加文件: {os.path.basename(file)}")
    
    def remove_selected_files(self):
        if self.is_processing:
            return
        
        selected = self.file_listbox.curselection()
        for index in reversed(selected):
            self.file_listbox.delete(index)
    
    def clear_file_list(self):
        if self.is_processing:
            return
        
        self.file_listbox.delete(0, tk.END)
    
    def select_output_file(self):
        if self.is_processing:
            return
        
        file_path = filedialog.asksaveasfilename(
            title="保存文件",
            defaultextension=".xlsx",
            filetypes=[
                ("Excel文件", "*.xlsx"),
                ("CSV文件", "*.csv"),
                ("所有文件", "*.*")
            ]
        )
        if file_path:
            self.output_path.set(file_path)
    
    def select_input_file(self):
        if self.is_processing:
            return
        
        convert_type = self.convert_type.get()
        
        if convert_type in ["excel_to_word", "excel_to_pdf", "excel_to_ppt"]:
            file_path = filedialog.askopenfilename(
                title="选择文件",
                filetypes=[
                    ("所有支持的文件", "*.xlsx *.xls *.csv"),
                    ("Excel文件", "*.xlsx *.xls"),
                    ("CSV文件", "*.csv")
                ]
            )
        elif convert_type in ["word_to_excel", "word_to_pdf"]:
            file_path = filedialog.askopenfilename(
                title="选择Word文件",
                filetypes=[("Word文件", "*.docx *.doc")]
            )
        elif convert_type == "ppt_to_pdf":
            file_path = filedialog.askopenfilename(
                title="选择PPT文件",
                filetypes=[("PPT文件", "*.pptx *.ppt")]
            )
        else:
            file_path = filedialog.askopenfilename(
                title="选择文件",
                filetypes=[("所有文件", "*.*")]
            )
        
        if file_path:
            self.convert_input.set(file_path)
    
    def select_convert_output(self):
        if self.is_processing:
            return
        
        convert_type = self.convert_type.get()
        
        if convert_type == "excel_to_word":
            file_path = filedialog.asksaveasfilename(
                title="保存Word文件",
                defaultextension=".docx",
                filetypes=[("Word文件", "*.docx")]
            )
        elif convert_type == "word_to_excel":
            file_path = filedialog.asksaveasfilename(
                title="保存文件",
                defaultextension=".xlsx",
                filetypes=[
                    ("Excel文件", "*.xlsx"),
                    ("CSV文件", "*.csv")
                ]
            )
        elif convert_type in ["excel_to_pdf", "word_to_pdf", "ppt_to_pdf"]:
            file_path = filedialog.asksaveasfilename(
                title="保存PDF文件",
                defaultextension=".pdf",
                filetypes=[("PDF文件", "*.pdf")]
            )
        elif convert_type == "excel_to_ppt":
            file_path = filedialog.asksaveasfilename(
                title="保存PPT文件",
                defaultextension=".pptx",
                filetypes=[("PPT文件", "*.pptx")]
            )
        else:
            file_path = filedialog.asksaveasfilename(
                title="保存文件",
                filetypes=[("所有文件", "*.*")]
            )
        
        if file_path:
            self.convert_output.set(file_path)
    
    def select_split_file(self):
        if self.is_processing:
            return
        
        file_path = filedialog.askopenfilename(
            title="选择文件",
            filetypes=[
                ("所有支持的文件", "*.xlsx *.xls *.csv"),
                ("Excel文件", "*.xlsx *.xls"),
                ("CSV文件", "*.csv")
            ]
        )
        if file_path:
            self.split_file_path.set(file_path)
    
    def select_split_output(self):
        if self.is_processing:
            return
        
        dir_path = filedialog.askdirectory(title="选择输出目录")
        if dir_path:
            self.split_output_dir.set(dir_path)
    
    def select_batch_dir(self):
        if self.is_processing:
            return
        
        dir_path = filedialog.askdirectory(title="选择文件夹")
        if dir_path:
            self.batch_dir.set(dir_path)
    
    def update_file_types(self):
        pass
    
    # 处理功能方法
    def start_merge(self):
        if self.is_processing:
            return
        
        files = list(self.file_listbox.get(0, tk.END))
        if not files:
            messagebox.showwarning("警告", "请先添加要合并的文件")
            return
        
        output_path = self.output_path.get()
        if not output_path:
            messagebox.showwarning("警告", "请指定输出文件路径")
            return
        
        thread = threading.Thread(target=self.merge_thread, args=(files, output_path))
        thread.daemon = True
        thread.start()
    
    def merge_thread(self, files, output_path):
        try:
            self.start_indeterminate("正在合并文件...")
            self.log("开始合并文件...")
            
            reference_columns = None
            
            for file in files:
                columns, has_header = self.get_columns_quick(file)
                if has_header:
                    reference_columns = columns
                    break
            
            if reference_columns is None:
                columns, _ = self.get_columns_quick(files[0])
                reference_columns = columns if columns else [f"列{i+1}" for i in range(10)]
            
            selected_columns = self.parse_column_selection(
                self.columns_by_name.get(), reference_columns
            )
            
            dfs = []
            total_files = len(files)
            
            for i, file in enumerate(files):
                self.check_cancel()
                
                if i == 0 or i == total_files - 1 or i % 5 == 0:
                    self.log(f"读取文件 {i+1}/{total_files}: {os.path.basename(file)}")
                
                self.update_progress((i / total_files) * 40, f"读取文件 {i+1}/{total_files}")
                
                df = self.read_excel_file(file)
                has_header = self.detect_header(df)
                
                if not has_header:
                    file_extension = os.path.splitext(file)[1].lower()
                    if file_extension == '.csv':
                        df = self.read_csv_without_header(file)
                    else:
                        df = self.read_excel_without_header(file)
                
                df = self.normalize_columns(df, reference_columns)
                
                if selected_columns:
                    available_columns = [col for col in selected_columns if col in df.columns]
                    df = df[available_columns]
                
                if self.remove_blank_rows_var.get() or self.remove_blank_cols_var.get() or self.remove_blank_cells_var.get():
                    df = self.remove_blank_data(
                        df,
                        remove_blank_rows=self.remove_blank_rows_var.get(),
                        remove_blank_cols=self.remove_blank_cols_var.get(),
                        remove_blank_cells=self.remove_blank_cells_var.get()
                    )
                
                if self.add_source_var.get():
                    df['数据来源'] = os.path.basename(file)
                
                dfs.append(df)
            
            self.check_cancel()
            self.update_progress(60, "正在合并数据...", force_update=True)
            
            merge_type = self.merge_type.get()
            if merge_type == "vertical":
                merged_df = pd.concat(dfs, ignore_index=True)
            else:
                merged_df = pd.concat(dfs, axis=1)
            
            self.update_progress(75, "正在处理数据...", force_update=True)
            
            if self.remove_blank_rows_var.get() or self.remove_blank_cols_var.get():
                merged_df = self.remove_blank_data(
                    merged_df,
                    remove_blank_rows=self.remove_blank_rows_var.get(),
                    remove_blank_cols=self.remove_blank_cols_var.get(),
                    remove_blank_cells=False
                )
            
            if self.remove_dup_var.get():
                self.check_cancel()
                merged_df = merged_df.drop_duplicates()
            
            self.check_cancel()
            self.update_progress(90, "正在保存结果...", force_update=True)
            
            if self.save_excel_file(merged_df, output_path):
                self.update_progress(100, "合并完成", force_update=True)
                self.log(f"合并完成！结果已保存到: {output_path}")
                
                self.stop_indeterminate("合并完成", success=True)
                messagebox.showinfo("成功", f"合并完成！\n输出文件: {output_path}")
            else:
                raise Exception("保存文件失败")
            
        except Exception as e:
            if "取消" in str(e):
                self.stop_indeterminate("已取消", success=False)
            else:
                self.stop_indeterminate("处理失败", success=False)
                self.log(f"合并失败: {str(e)}")
                messagebox.showerror("错误", f"合并失败: {str(e)}")
    
    def start_convert(self):
        if self.is_processing:
            return
        
        input_path = self.convert_input.get()
        output_path = self.convert_output.get()
        
        if not input_path:
            messagebox.showwarning("警告", "请选择输入文件")
            return
        if not output_path:
            messagebox.showwarning("警告", "请指定输出文件路径")
            return
        
        thread = threading.Thread(target=self.convert_thread, args=(input_path, output_path))
        thread.daemon = True
        thread.start()
    
    def convert_thread(self, input_path, output_path):
        try:
            convert_type = self.convert_type.get()
            self.start_indeterminate("正在转换...")
            
            if convert_type == "excel_to_word":
                self.convert_excel_to_word(input_path, output_path)
            elif convert_type == "word_to_excel":
                self.convert_word_to_excel(input_path, output_path)
            elif convert_type == "excel_to_pdf":
                self.convert_excel_to_pdf(input_path, output_path)
            elif convert_type == "excel_to_ppt":
                self.convert_excel_to_ppt(input_path, output_path)
            elif convert_type == "word_to_pdf":
                self.convert_word_to_pdf(input_path, output_path)
            elif convert_type == "ppt_to_pdf":
                self.convert_ppt_to_pdf(input_path, output_path)
            
            self.update_progress(100, "转换完成", force_update=True)
            self.stop_indeterminate("转换完成", success=True)
            messagebox.showinfo("成功", f"转换完成！\n输出文件: {output_path}")
            
        except Exception as e:
            if "取消" in str(e):
                self.stop_indeterminate("已取消", success=False)
            else:
                self.stop_indeterminate("转换失败", success=False)
                self.log(f"转换失败: {str(e)}")
                messagebox.showerror("错误", f"转换失败: {str(e)}")
    
    def convert_excel_to_word(self, input_path, output_path):
        """Excel转Word"""
        from docx import Document
        
        self.update_progress(20, "正在读取文件...", force_update=True)
        df = self.read_excel_file(input_path)
        
        if not self.detect_header(df):
            file_extension = os.path.splitext(input_path)[1].lower()
            if file_extension == '.csv':
                df = self.read_csv_without_header(input_path)
            else:
                df = self.read_excel_without_header(input_path)
        
        df = self.remove_blank_data(df)
        df = self.convert_dtypes_after_processing(df)
        
        self.check_cancel()
        self.update_progress(50, "正在创建Word文档...", force_update=True)
        doc = Document()
        doc.add_heading('数据转换结果', level=1)
        
        table = doc.add_table(rows=1, cols=len(df.columns))
        table.style = 'Light Grid Accent 1'
        
        header_cells = table.rows[0].cells
        for i, col in enumerate(df.columns):
            header_cells[i].text = str(col)
        
        total_rows = len(df)
        batch_size = max(1, total_rows // 10)
        
        for idx, (_, row) in enumerate(df.iterrows()):
            self.check_cancel()
            
            row_cells = table.add_row().cells
            for i, value in enumerate(row):
                if pd.isna(value):
                    row_cells[i].text = ''
                elif isinstance(value, (int, float, np.int64, np.float64)):
                    if isinstance(value, float) and value.is_integer():
                        row_cells[i].text = str(int(value))
                    else:
                        row_cells[i].text = str(value)
                else:
                    row_cells[i].text = str(value)
            
            if idx % batch_size == 0:
                progress = 50 + (idx / total_rows) * 40
                self.update_progress(progress)
        
        self.update_progress(90, "正在保存Word文件...", force_update=True)
        doc.save(output_path)
        self.log(f"Excel转Word完成: {output_path}")
    
    def convert_word_to_excel(self, input_path, output_path):
        """Word转Excel"""
        from docx import Document
        
        self.update_progress(30, "正在读取Word文件...", force_update=True)
        doc = Document(input_path)
        
        if not doc.tables:
            raise Exception("Word文档中没有表格")
        
        self.check_cancel()
        self.update_progress(60, "正在提取表格数据...", force_update=True)
        table = doc.tables[0]
        data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            data.append(row_data)
        
        self.update_progress(80, "正在转换为Excel...", force_update=True)
        df = pd.DataFrame(data[1:], columns=data[0])
        
        df = self.remove_blank_data(df)
        
        if self.save_excel_file(df, output_path):
            self.log(f"Word转Excel完成: {output_path}")
        else:
            raise Exception("保存文件失败")
    
    def convert_excel_to_pdf(self, input_path, output_path):
        """Excel转PDF"""
        try:
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            
            self.update_progress(20, "正在读取Excel文件...", force_update=True)
            df = self.read_excel_file(input_path)
            
            if not self.detect_header(df):
                file_extension = os.path.splitext(input_path)[1].lower()
                if file_extension == '.csv':
                    df = self.read_csv_without_header(input_path)
                else:
                    df = self.read_excel_without_header(input_path)
            
            df = self.remove_blank_data(df)
            df = self.convert_dtypes_after_processing(df)
            
            self.check_cancel()
            self.update_progress(50, "正在创建PDF...", force_update=True)
            
            doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
            elements = []
            
            table_data = [df.columns.tolist()] + df.values.tolist()
            table = Table(table_data)
            
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
            ])
            table.setStyle(style)
            
            elements.append(table)
            
            self.check_cancel()
            self.update_progress(80, "正在保存PDF...", force_update=True)
            
            doc.build(elements)
            self.log(f"Excel转PDF完成: {output_path}")
            
        except ImportError:
            raise Exception("缺少必要的库: reportlab")
    
    def convert_excel_to_ppt(self, input_path, output_path):
        """Excel转PPT"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            self.update_progress(20, "正在读取Excel文件...", force_update=True)
            df = self.read_excel_file(input_path)
            
            if not self.detect_header(df):
                file_extension = os.path.splitext(input_path)[1].lower()
                if file_extension == '.csv':
                    df = self.read_csv_without_header(input_path)
                else:
                    df = self.read_excel_without_header(input_path)
            
            df = self.remove_blank_data(df)
            df = self.convert_dtypes_after_processing(df)
            
            self.check_cancel()
            self.update_progress(50, "正在创建PPT...", force_update=True)
            
            prs = Presentation()
            
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "数据展示"
            subtitle.text = f"数据行数: {len(df)}"
            
            table_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(table_slide_layout)
            title = slide.shapes.title
            title.text = "数据详情"
            
            rows, cols = min(len(df) + 1, 20), min(len(df.columns), 8)
            table_shape = slide.shapes.add_table(rows, cols, 
                                                 Inches(0.5), Inches(1.5),
                                                 Inches(9), Inches(5))
            table = table_shape.table
            
            for i, col in enumerate(df.columns[:cols]):
                table.cell(0, i).text = str(col)
            
            for i in range(1, rows):
                for j in range(cols):
                    value = df.iloc[i-1, j]
                    if pd.isna(value):
                        table.cell(i, j).text = ''
                    else:
                        table.cell(i, j).text = str(value)
            
            self.check_cancel()
            self.update_progress(80, "正在保存PPT...", force_update=True)
            
            prs.save(output_path)
            self.log(f"Excel转PPT完成: {output_path}")
            
        except ImportError:
            raise Exception("缺少必要的库: python-pptx")
    
    def convert_word_to_pdf(self, input_path, output_path):
        """Word转PDF"""
        try:
            try:
                from docx2pdf import convert
                self.update_progress(50, "正在转换...", force_update=True)
                convert(input_path, output_path)
                self.log(f"Word转PDF完成: {output_path}")
                return
            except ImportError:
                pass
            
            from docx import Document
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            self.update_progress(30, "正在读取Word文件...", force_update=True)
            doc = Document(input_path)
            
            self.check_cancel()
            self.update_progress(60, "正在创建PDF...", force_update=True)
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            elements = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    elements.append(Paragraph(para.text, styles['Normal']))
                    elements.append(Spacer(1, 12))
            
            self.update_progress(80, "正在保存PDF...", force_update=True)
            pdf_doc.build(elements)
            self.log(f"Word转PDF完成: {output_path}")
            
        except ImportError:
            raise Exception("缺少必要的库: docx2pdf 或 reportlab")
    
    def convert_ppt_to_pdf(self, input_path, output_path):
        """PPT转PDF"""
        try:
            try:
                import win32com.client
                
                self.update_progress(50, "正在转换...", force_update=True)
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(input_path)
                presentation.SaveAs(output_path, 32)
                presentation.Close()
                powerpoint.Quit()
                
                self.log(f"PPT转PDF完成: {output_path}")
                return
                
            except ImportError:
                pass
            
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            self.update_progress(30, "正在读取PPT文件...", force_update=True)
            prs = Presentation(input_path)
            
            self.check_cancel()
            self.update_progress(60, "正在创建PDF...", force_update=True)
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
            styles = getSampleStyleSheet()
            elements = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                elements.append(Paragraph(para.text, styles['Normal']))
                                elements.append(Spacer(1, 12))
                
                elements.append(Spacer(1, 30))
            
            self.update_progress(80, "正在保存PDF...", force_update=True)
            pdf_doc.build(elements)
            self.log(f"PPT转PDF完成: {output_path}")
            
        except ImportError:
            raise Exception("缺少必要的库: pywin32 或 python-pptx")
    
    def start_split(self):
        if self.is_processing:
            return
        
        file_path = self.split_file_path.get()
        if not file_path:
            messagebox.showwarning("警告", "请选择要拆分的文件")
            return
        
        thread = threading.Thread(target=self.split_thread, args=(file_path,))
        thread.daemon = True
        thread.start()
    
    def split_thread(self, file_path):
        try:
            self.start_indeterminate("正在拆分数据...")
            
            self.update_progress(10, "正在读取文件...", force_update=True)
            df = self.read_excel_file(file_path)
            
            if not self.detect_header(df):
                file_extension = os.path.splitext(file_path)[1].lower()
                if file_extension == '.csv':
                    df = self.read_csv_without_header(file_path)
                else:
                    df = self.read_excel_without_header(file_path)
            
            self.log("正在去除空白数据...")
            df = self.remove_blank_data(df)
            
            output_dir = self.split_output_dir.get()
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            split_method = self.split_method.get()
            
            if split_method == "by_columns":
                self.split_by_columns(df, output_dir)
            elif split_method == "by_column_position":
                self.split_by_column_position(df, output_dir)
            elif split_method == "by_rows":
                self.split_by_rows(df, output_dir)
            elif split_method == "by_specific_rows":
                self.split_by_specific_rows(df, output_dir)
            
            self.update_progress(100, "拆分完成", force_update=True)
            self.log("拆分完成！")
            
            self.stop_indeterminate("拆分完成", success=True)
            messagebox.showinfo("成功", f"拆分完成！\n输出目录: {output_dir}")
            
        except Exception as e:
            if "取消" in str(e):
                self.stop_indeterminate("已取消", success=False)
            else:
                self.stop_indeterminate("拆分失败", success=False)
                self.log(f"拆分失败: {str(e)}")
                messagebox.showerror("错误", f"拆分失败: {str(e)}")
    
    def split_by_columns(self, df, output_dir):
        """按列值拆分"""
        columns_str = self.split_columns.get().strip()
        if not columns_str:
            raise Exception("请输入拆分依据列")
        
        split_columns = [c.strip() for c in columns_str.split(',') if c.strip()]
        
        missing_columns = [col for col in split_columns if col not in df.columns]
        if missing_columns:
            raise Exception(f"列不存在: {missing_columns}")
        
        groups = df.groupby(split_columns)
        total_groups = len(groups)
        
        for i, (name, group) in enumerate(groups):
            self.check_cancel()
            
            if isinstance(name, tuple):
                safe_name = '_'.join([re.sub(r'[\\/*?:"<>|]', '_', str(n)) for n in name])
            else:
                safe_name = re.sub(r'[\\/*?:"<>|]', '_', str(name))
            
            output_path = os.path.join(output_dir, f"{safe_name}.xlsx")
            
            if self.save_excel_file(group, output_path):
                if i % max(1, total_groups // 20) == 0:
                    progress = 20 + ((i + 1) / total_groups) * 70
                    self.update_progress(progress)
        
        self.log(f"按列拆分完成！共生成 {total_groups} 个文件")
    
    def split_by_column_position(self, df, output_dir):
        """按列位置拆分"""
        pos_type = self.column_position_type.get()
        
        if pos_type == "before_after":
            try:
                split_pos = int(self.split_position.get().strip())
                if split_pos < 1 or split_pos >= len(df.columns):
                    raise ValueError(f"拆分位置必须在1到{len(df.columns)-1}之间")
            except ValueError as e:
                raise Exception(f"无效的拆分位置: {e}")
            
            df_before = df.iloc[:, :split_pos]
            output_path_before = os.path.join(output_dir, f"列1-{split_pos}.xlsx")
            
            self.check_cancel()
            if self.save_excel_file(df_before, output_path_before):
                self.update_progress(50, force_update=True)
                self.log(f"已保存: 列1-{split_pos}.xlsx")
            
            df_after = df.iloc[:, split_pos:]
            output_path_after = os.path.join(output_dir, f"列{split_pos+1}-{len(df.columns)}.xlsx")
            
            self.check_cancel()
            if self.save_excel_file(df_after, output_path_after):
                self.update_progress(100, force_update=True)
                self.log(f"已保存: 列{split_pos+1}-{len(df.columns)}.xlsx")
            
            self.log("按列位置前后拆分完成！共生成 2 个文件")
            
        elif pos_type == "specific_columns":
            columns_str = self.specific_column_numbers.get().strip()
            if not columns_str:
                raise Exception("请输入指定列号")
            
            try:
                column_indices = []
                for part in columns_str.split(','):
                    part = part.strip()
                    if '-' in part:
                        start, end = part.split('-')
                        column_indices.extend(range(int(start), int(end) + 1))
                    else:
                        column_indices.append(int(part))
                
                column_indices = sorted(set(column_indices))
                
                valid_indices = [i for i in column_indices if 1 <= i <= len(df.columns)]
                
                if not valid_indices:
                    raise Exception("没有有效的列号")
                
            except ValueError:
                raise Exception("无效的列号格式")
            
            df_specific = df.iloc[:, [i-1 for i in valid_indices]]
            
            output_path = os.path.join(output_dir, f"指定列_{'_'.join(map(str, valid_indices))}.xlsx")
            
            self.check_cancel()
            if self.save_excel_file(df_specific, output_path):
                self.update_progress(100, force_update=True)
                self.log(f"已保存: 指定列_{'_'.join(map(str, valid_indices))}.xlsx")
            
            self.log("按指定列提取完成！共生成 1 个文件")
    
    def split_by_rows(self, df, output_dir):
        """按行数拆分"""
        try:
            rows_per_file = int(self.rows_per_file.get())
            if rows_per_file <= 0:
                raise ValueError("每个文件行数必须大于0")
        except ValueError as e:
            raise Exception(f"无效的行数: {e}")
        
        total_rows = len(df)
        total_files = (total_rows + rows_per_file - 1) // rows_per_file
        
        for i in range(total_files):
            self.check_cancel()
            
            start_idx = i * rows_per_file
            end_idx = min((i + 1) * rows_per_file, total_rows)
            
            chunk = df.iloc[start_idx:end_idx]
            
            output_path = os.path.join(output_dir, f"part_{i+1:03d}.xlsx")
            
            if self.save_excel_file(chunk, output_path):
                progress = 20 + ((i + 1) / total_files) * 70
                self.update_progress(progress)
        
        self.log(f"按行数拆分完成！共生成 {total_files} 个文件")
    
    def split_by_specific_rows(self, df, output_dir):
        """按特定行拆分"""
        rows_str = self.specific_rows.get().strip()
        if not rows_str:
            raise Exception("请输入拆分行号")
        
        try:
            split_rows = [int(r.strip()) for r in rows_str.split(',') if r.strip()]
            split_rows.sort()
        except ValueError:
            raise Exception("无效的行号格式")
        
        if not split_rows:
            raise Exception("请输入有效的行号")
        
        all_split_points = [0] + split_rows + [len(df)]
        total_files = len(split_rows) + 1
        
        for i in range(total_files):
            self.check_cancel()
            
            start_idx = all_split_points[i]
            end_idx = all_split_points[i + 1]
            
            chunk = df.iloc[start_idx:end_idx]
            
            output_path = os.path.join(output_dir, f"part_{i+1:03d}.xlsx")
            
            if self.save_excel_file(chunk, output_path):
                progress = 20 + ((i + 1) / total_files) * 70
                self.update_progress(progress)
        
        self.log(f"按特定行拆分完成！共生成 {total_files} 个文件")
    
    def start_batch(self):
        if self.is_processing:
            return
        
        directory = self.batch_dir.get()
        if not directory:
            messagebox.showwarning("警告", "请选择要处理的文件夹")
            return
        
        thread = threading.Thread(target=self.batch_thread, args=(directory,))
        thread.daemon = True
        thread.start()
    
    def batch_thread(self, directory):
        try:
            self.start_indeterminate("正在批量处理...")
            
            excel_files = [f for f in os.listdir(directory) 
                          if f.endswith(('.xlsx', '.xls', '.csv'))]
            
            if not excel_files:
                self.stop_indeterminate("未找到文件", success=False)
                messagebox.showwarning("警告", "目录中没有支持的文件")
                return
            
            output_dir = os.path.join(directory, 'processed')
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            operation = self.batch_operation.get()
            total_files = len(excel_files)
            
            for i, file_name in enumerate(excel_files):
                self.check_cancel()
                file_path = os.path.join(directory, file_name)
                
                if i % 5 == 0 or i == total_files - 1:
                    self.log(f"处理文件 {i+1}/{total_files}: {file_name}")
                
                progress = (i / total_files) * 100
                self.update_progress(progress)
                
                df = self.read_excel_file(file_path)
                
                if not self.detect_header(df):
                    file_extension = os.path.splitext(file_path)[1].lower()
                    if file_extension == '.csv':
                        df = self.read_csv_without_header(file_path)
                    else:
                        df = self.read_excel_without_header(file_path)
                
                if operation == "clean":
                    df = self.remove_blank_data(df)
                    df = df.drop_duplicates()
                elif operation == "remove_empty":
                    df = df.dropna()
                
                output_path = os.path.join(output_dir, f"processed_{file_name}")
                self.save_excel_file(df, output_path)
            
            self.update_progress(100, "批量处理完成", force_update=True)
            self.log(f"批量处理完成！处理了 {total_files} 个文件")
            
            self.stop_indeterminate("批量处理完成", success=True)
            messagebox.showinfo("成功", f"批量处理完成！\n处理了 {total_files} 个文件")
            
        except Exception as e:
            if "取消" in str(e):
                self.stop_indeterminate("已取消", success=False)
            else:
                self.stop_indeterminate("处理失败", success=False)
                self.log(f"批量处理失败: {str(e)}")
                messagebox.showerror("错误", f"批量处理失败: {str(e)}")

def main():
    root = tk.Tk()
    app = DataProcessingGUI(root)
    root.mainloop()

if __name__ == "__main__":
    main()
